import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create the presentation object
    prs = Presentation()

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Rape Cases in India: 2000–2026 Context and Data"
    subtitle.text = "Year-Wise Statistics, Political Context, and Judicial Release Information\nSource: NCRB Reports"

    # 2. Context & Data Warning Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Context & Data Availability"
    tf = body_shape.text_frame
    tf.text = "Data Source: National Crime Records Bureau (NCRB) - 'Crime in India' reports."
    
    p = tf.add_paragraph()
    p.text = "The numbers reflect registered cases under Section 376 IPC (and later POCSO)."
    
    p = tf.add_paragraph()
    p.text = "An increase in numbers post-2012 points to better reporting and broader legal definitions (following the Nirbhaya case), rather than purely an increase in crime rate."
    
    p = tf.add_paragraph()
    p.text = "Data for 2024–2026: Official consolidated NCRB reports trail by 1-2 years, so exact figures for these years are pending compilation."

    # 3. Year-wise data Slides
    # We will split data into two slides to fit nicely
    data = [
        ("2000", "16,496", "NDA (BJP-led)", "Atal Bihari Vajpayee"),
        ("2001", "16,075", "NDA (BJP-led)", "Atal Bihari Vajpayee"),
        ("2002", "16,373", "NDA (BJP-led)", "Atal Bihari Vajpayee"),
        ("2003", "15,847", "NDA (BJP-led)", "Atal Bihari Vajpayee"),
        ("2004", "18,233", "NDA / UPA (INC-led)", "A.B. Vajpayee / Manmohan Singh"),
        ("2005", "18,359", "UPA (INC-led)", "Manmohan Singh"),
        ("2006", "19,348", "UPA (INC-led)", "Manmohan Singh"),
        ("2007", "20,737", "UPA (INC-led)", "Manmohan Singh"),
        ("2008", "21,467", "UPA (INC-led)", "Manmohan Singh"),
        ("2009", "21,397", "UPA (INC-led)", "Manmohan Singh"),
        ("2010", "22,172", "UPA (INC-led)", "Manmohan Singh"),
        ("2011", "24,206", "UPA (INC-led)", "Manmohan Singh"),
        ("2012", "24,923", "UPA (INC-led)", "Manmohan Singh"),
        ("2013", "33,707*", "UPA (INC-led)", "Manmohan Singh"),
        ("2014", "36,735", "UPA / NDA (BJP-led)", "Manmohan / Narendra Modi"),
        ("2015", "34,651", "NDA (BJP-led)", "Narendra Modi"),
        ("2016", "38,947", "NDA (BJP-led)", "Narendra Modi"),
        ("2017", "32,559", "NDA (BJP-led)", "Narendra Modi"),
        ("2018", "33,356", "NDA (BJP-led)", "Narendra Modi"),
        ("2019", "32,033", "NDA (BJP-led)", "Narendra Modi"),
        ("2020", "28,046", "NDA (BJP-led)", "Narendra Modi"),
        ("2021", "31,677", "NDA (BJP-led)", "Narendra Modi"),
        ("2022", "31,516", "NDA (BJP-led)", "Narendra Modi"),
        ("2023", "29,670", "NDA (BJP-led)", "Narendra Modi"),
        ("2024", "Pending", "NDA (BJP-led)", "Narendra Modi"),
        ("2025", "Pending", "NDA (BJP-led)", "Narendra Modi"),
        ("2026", "Pending", "NDA (BJP-led)", "Narendra Modi"),
    ]

    # Helper function to create table slides
    def create_table_slide(title_text, row_data):
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
        slide.shapes.title.text = title_text
        
        rows = len(row_data) + 1
        cols = 4
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(0.4 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(1.0)
        table.columns[1].width = Inches(1.8)
        table.columns[2].width = Inches(2.8)
        table.columns[3].width = Inches(3.4)

        headers = ["Year", "Cases Registered", "Ruling Party", "Prime Minister"]
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            # basic styling
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                
        for row_idx, row_item in enumerate(row_data):
            for col_idx, val in enumerate(row_item):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)

    # Split data to fit neatly on slides
    create_table_slide("NCRB Data (2000 \u2013 2013)", data[:14])
    create_table_slide("NCRB Data (2014 \u2013 2026)", data[14:])

    # 4. Convicted Release context
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Release of Convicted Accused: Overview"
    tf = body_shape.text_frame
    tf.text = "NCRB does not publish national aggregated statistics specifically on convicted rapists released."
    
    p = tf.add_paragraph()
    p.text = "Overall Conviction Rate: Remains low (22% - 28%). Acquittals form the majority due to hostile witnesses, lack of evidence, and delays."
    
    p = tf.add_paragraph()
    p.text = "Prison administration is a 'State Subject'. Premature release decisions are made by State Sentence Review Boards (SSRBs)."

    # 5. Reasons slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Reasons for Release"
    tf = body_shape.text_frame
    tf.text = "1. Remission of Sentence (Premature Release):"
    
    p = tf.add_paragraph()
    p.text = "Provided for 'good behavior', finishing 14 years jail term. Ex: 11 Bilkis Bano convicts were originally released by Gujarat govt via its 1992 policy (this was later struck down by the Supreme Court as illegal in 2024)."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "2. Parole / Furlough:"
    p = tf.add_paragraph()
    p.text = "Temporary release for medical/family emergencies or societal ties. Ex: High profile cases like Gurmeet Ram Rahim getting recurrent parole, which sometimes sparks political debate."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Completion of Standard Term:"
    p = tf.add_paragraph()
    p.text = "For cases without life imprisonment, standard terms are finished routinely."
    p.level = 1

    file_name = "India_Rape_Cases_Data_2000_2026.pptx"
    prs.save(file_name)
    print(f"Presentation saved successfully as {file_name}")

if __name__ == '__main__':
    create_presentation()
